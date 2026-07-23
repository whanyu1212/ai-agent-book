"""
Document processing tools for PDF, DOCX, PPTX, CSV, TXT.
Based on AWorld MCP server implementation.
"""
import json
import logging
import traceback
from pathlib import Path
from typing import Union, Dict, Any

import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2
from dotenv import load_dotenv
from mcp.types import TextContent

from base import ActionResponse, validate_file_path


load_dotenv()


async def extract_pdf_text(
    file_path: str,
    page_range: str | None = None
) -> Union[str, TextContent]:
    """
    Extract text from PDF file.
    
    Args:
        file_path: Path to PDF file
        page_range: Optional page range (e.g., "1-5" or "1,3,5")
        
    Returns:
        TextContent with extracted text
    """
    try:
        path = validate_file_path(file_path)
        
        logging.info(f"📄 Extracting PDF: {path}")
        
        with open(path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            total_pages = len(reader.pages)
            
            # Parse page range
            if page_range:
                pages_to_extract = parse_page_range(page_range, total_pages)
            else:
                pages_to_extract = range(total_pages)
            
            # Extract text
            text_parts = []
            for page_num in pages_to_extract:
                if page_num < total_pages:
                    page = reader.pages[page_num]
                    text = page.extract_text()
                    text_parts.append(f"--- Page {page_num + 1} ---\n{text}\n")
            
            full_text = "\n".join(text_parts)
            
            result = {
                "file_name": path.name,
                "file_type": "pdf",
                "total_pages": total_pages,
                "pages_extracted": len(pages_to_extract),
                "text": full_text[:50000],  # Limit to 50k chars
                "text_length": len(full_text),
                "truncated": len(full_text) > 50000
            }
            
            logging.info(f"✅ Extracted {len(pages_to_extract)} pages from PDF")
            
            action_response = ActionResponse(
                success=True,
                message=result,
                metadata={"file_path": str(path), "pages": total_pages}
            )
            
            return TextContent(
                type="text",
                text=json.dumps(action_response.model_dump())
            )
            
    except Exception as e:
        error_msg = f"PDF extraction failed: {str(e)}"
        logging.error(f"PDF error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "pdf_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def extract_docx_content(
    file_path: str
) -> Union[str, TextContent]:
    """
    Extract content from DOCX file.
    
    Args:
        file_path: Path to DOCX file
        
    Returns:
        TextContent with extracted content
    """
    try:
        path = validate_file_path(file_path)
        
        logging.info(f"📄 Extracting DOCX: {path}")
        
        doc = Document(path)
        
        # Extract paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Extract tables
        tables_data = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables_data.append(table_data)
        
        full_text = "\n\n".join(paragraphs)
        
        result = {
            "file_name": path.name,
            "file_type": "docx",
            "paragraphs": len(paragraphs),
            "tables": len(tables_data),
            "text": full_text[:50000],
            "text_length": len(full_text),
            "truncated": len(full_text) > 50000,
            "tables_data": tables_data if tables_data else []
        }
        
        logging.info(f"✅ Extracted DOCX: {len(paragraphs)} paragraphs, {len(tables_data)} tables")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"DOCX extraction failed: {str(e)}"
        logging.error(f"DOCX error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "docx_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def extract_pptx_content(
    file_path: str
) -> Union[str, TextContent]:
    """
    Extract content from PPTX file.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        TextContent with extracted content
    """
    try:
        path = validate_file_path(file_path)
        
        logging.info(f"📊 Extracting PPTX: {path}")
        
        prs = Presentation(path)
        
        slides_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                slides_content.append({
                    "slide_number": slide_num,
                    "text": "\n".join(slide_text)
                })
        
        full_text = "\n\n".join([f"=== Slide {s['slide_number']} ===\n{s['text']}" for s in slides_content])
        
        result = {
            "file_name": path.name,
            "file_type": "pptx",
            "total_slides": len(prs.slides),
            "slides_with_content": len(slides_content),
            "text": full_text[:50000],
            "text_length": len(full_text),
            "truncated": len(full_text) > 50000,
            "slides": slides_content
        }
        
        logging.info(f"✅ Extracted PPTX: {len(prs.slides)} slides")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"PPTX extraction failed: {str(e)}"
        logging.error(f"PPTX error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "pptx_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def extract_csv_content(
    file_path: str,
    max_rows: int = 1000
) -> Union[str, TextContent]:
    """
    Extract and parse CSV file content.
    
    Args:
        file_path: Path to CSV file
        max_rows: Maximum rows to read
        
    Returns:
        TextContent with parsed CSV data
    """
    try:
        path = validate_file_path(file_path)
        
        logging.info(f"📊 Parsing CSV: {path}")
        
        # Read CSV with pandas
        df = pd.read_csv(path, nrows=max_rows)
        
        result = {
            "file_name": path.name,
            "file_type": "csv",
            "rows": len(df),
            "columns": len(df.columns),
            "column_names": df.columns.tolist(),
            "data": df.to_dict(orient="records"),
            "preview": df.head(10).to_string(),
            "truncated": len(df) == max_rows
        }
        
        logging.info(f"✅ Parsed CSV: {len(df)} rows, {len(df.columns)} columns")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"CSV parsing failed: {str(e)}"
        logging.error(f"CSV error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "csv_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


def parse_page_range(page_range: str, total_pages: int) -> list[int]:
    """
    Parse page range string into list of page numbers.
    
    Args:
        page_range: String like "1-5" or "1,3,5" or "1-3,7,9-11"
        total_pages: Total number of pages
        
    Returns:
        List of page numbers (0-indexed)
    """
    pages = []
    
    for part in page_range.split(","):
        part = part.strip()
        if not part:
            # Trailing/duplicate commas (e.g. "1,3," or "1,,3") are common in
            # LLM tool args; skip empty segments instead of int("").
            continue
        if "-" in part:
            bounds = part.split("-")
            if len(bounds) != 2 or not bounds[0] or not bounds[1]:
                raise ValueError(f"Invalid page range segment: {part!r}")
            start, end = int(bounds[0]), int(bounds[1])
            # Clamp both ends: the caller's guard is `page_num < total_pages`,
            # which a negative index passes, and reader.pages[-1] is the LAST
            # page -- so an unclamped start silently returns the wrong page.
            pages.extend(range(max(0, start - 1), min(end, total_pages)))
        else:
            page_num = int(part) - 1
            if 0 <= page_num < total_pages:
                pages.append(page_num)
    
    return sorted(set(pages))
