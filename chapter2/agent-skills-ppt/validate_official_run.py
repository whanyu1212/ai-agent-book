#!/usr/bin/env python3
"""Fail-closed validator for the real agent-runtime + official PPTX Skill run.

Supports both accepted runtimes: Claude Code (claude_stream.jsonl) and Kimi
Code CLI (kimi_stream.jsonl). The artifact gates are identical; only the
stream parsing and run-success gate are runtime-specific.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def iter_values(value: Any):
    yield value
    if isinstance(value, dict):
        for item in value.values():
            yield from iter_values(item)
    elif isinstance(value, list):
        for item in value:
            yield from iter_values(item)


def parse_stream(path: Path) -> tuple[list[dict], str]:
    events = []
    raw = path.read_text(encoding="utf-8")
    for line in raw.splitlines():
        if not line.strip():
            continue
        try:
            events.append(json.loads(line))
        except json.JSONDecodeError:
            continue
    return events, raw


def collect_tool_calls(events: list[dict]) -> list[dict]:
    """Normalize tool calls across runtime stream formats.

    Claude Code emits ``assistant`` events whose message content contains
    ``tool_use`` blocks; Kimi Code CLI emits ``assistant`` events with an
    OpenAI-style ``tool_calls`` list.
    """
    calls = []
    for event in events:
        if event.get("role") == "assistant" and event.get("tool_calls"):
            for call in event["tool_calls"]:
                function = call.get("function", {})
                calls.append({
                    "name": function.get("name", ""),
                    "arguments": function.get("arguments", ""),
                })
        message = event.get("message")
        if isinstance(message, dict):
            for block in message.get("content") or []:
                if isinstance(block, dict) and block.get("type") == "tool_use":
                    calls.append({
                        "name": block.get("name", ""),
                        "arguments": json.dumps(block.get("input", {}), ensure_ascii=False),
                    })
    return calls


def collect_result_metadata(events: list[dict]) -> dict:
    result_events = [event for event in events if event.get("type") == "result"]
    if not result_events:
        return {}
    result = result_events[-1]
    return {
        key: result.get(key)
        for key in (
            "subtype", "is_error", "duration_ms", "duration_api_ms", "num_turns",
            "result", "total_cost_usd", "usage", "modelUsage", "session_id",
        )
    }


def collect_kimi_metadata(events: list[dict], run_dir: Path) -> dict:
    exit_path = run_dir / "kimi_exit.json"
    exit_info = json.loads(exit_path.read_text(encoding="utf-8")) if exit_path.exists() else {}
    tool_calls = collect_tool_calls(events)
    assistant_messages = [
        event.get("content", "")
        for event in events
        if event.get("role") == "assistant" and event.get("content")
    ]
    session_ids = [
        event.get("session_id")
        for event in events
        if event.get("role") == "meta" and event.get("session_id")
    ]
    runtime_info_path = run_dir / "runtime.json"
    runtime_info = (
        json.loads(runtime_info_path.read_text(encoding="utf-8"))
        if runtime_info_path.exists() else {}
    )
    return {
        "return_code": exit_info.get("return_code"),
        "model_alias": runtime_info.get("model_alias"),
        "num_assistant_messages": len(assistant_messages),
        "num_tool_calls": len(tool_calls),
        "tool_names": sorted({call["name"] for call in tool_calls if call["name"]}),
        "session_id": session_ids[-1] if session_ids else None,
        "final_response": assistant_messages[-1] if assistant_messages else None,
    }


def extract_slide_text(prs: Presentation) -> str:
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                chunks.append(shape.text)
    return "\n".join(chunks)


def validate(run_dir: Path) -> dict:
    run_dir = run_dir.resolve()
    protocol_path = run_dir / "experiment_protocol.json"
    protocol = json.loads(protocol_path.read_text(encoding="utf-8"))
    workspace = run_dir / "workspace"
    kimi_stream_path = run_dir / "kimi_stream.jsonl"
    runtime = "kimi" if kimi_stream_path.exists() else "claude"
    stream_path = kimi_stream_path if runtime == "kimi" else run_dir / "claude_stream.jsonl"
    paper_path = workspace / "attention-is-all-you-need.pdf"
    pptx_path = workspace / "output" / "attention-is-all-you-need.pptx"
    visual_dir = workspace / "source_visuals"
    visual_manifest_path = visual_dir / "manifest.json"
    events, raw_stream = parse_stream(stream_path)
    tool_calls = collect_tool_calls(events)
    tool_args_text = "\n".join(
        call["arguments"] for call in tool_calls if isinstance(call["arguments"], str)
    ).lower()

    prs = Presentation(str(pptx_path)) if pptx_path.exists() else None
    slide_count = len(prs.slides) if prs else 0
    slide_text = extract_slide_text(prs).lower() if prs else ""
    zip_valid = False
    embedded_media: dict[str, str] = {}
    if pptx_path.exists():
        with zipfile.ZipFile(pptx_path) as archive:
            zip_valid = archive.testzip() is None
            for name in archive.namelist():
                if name.startswith("ppt/media/") and not name.endswith("/"):
                    embedded_media[name] = hashlib.sha256(archive.read(name)).hexdigest()

    visual_manifest = (
        json.loads(visual_manifest_path.read_text(encoding="utf-8"))
        if visual_manifest_path.exists() else []
    )
    if isinstance(visual_manifest, dict):
        visual_manifest = visual_manifest.get("visuals", [])
    source_visuals = []
    for item in visual_manifest if isinstance(visual_manifest, list) else []:
        filename = item.get("file") if isinstance(item, dict) else None
        path = visual_dir / filename if filename else None
        if path and path.is_file():
            source_visuals.append({
                "file": filename,
                "sha256": sha256(path),
                "embedded": sha256(path) in embedded_media.values(),
                "page": item.get("page"),
                "label": item.get("label"),
                "caption": item.get("caption"),
            })

    thumbnail_candidates = sorted((workspace / "output").glob("*thumbnail*.jpg"))
    if not thumbnail_candidates:
        thumbnail_candidates = sorted((workspace / "output").glob("*thumbnails*.jpg"))
    values_as_text = "\n".join(str(value) for event in events for value in iter_values(event))
    lower_evidence = (raw_stream + "\n" + values_as_text).lower()
    section_checks = {
        "title": "attention is all you need" in slide_text,
        "background": any(term in slide_text for term in ("background", "motivation", "problem")),
        "method": "transformer" in slide_text and any(term in slide_text for term in ("architecture", "method")),
        "results": any(term in slide_text for term in ("result", "bleu", "translation")),
        "conclusion": any(term in slide_text for term in ("conclusion", "takeaway", "summary")),
    }
    # Progressive-disclosure evidence must come from what the agent actually
    # did (tool calls and tool results), not merely from the prompt text.
    if runtime == "kimi":
        skill_calls = [
            call for call in tool_calls
            if call["name"] == "Skill" and '"pptx"' in call["arguments"].lower().replace(" ", "")
        ]
        pptx_skill_invoked = bool(skill_calls)
        # Kimi's Skill tool loads the complete SKILL.md inline by construction;
        # require the load to have succeeded in the recorded tool result.
        skill_md_loaded = pptx_skill_invoked and 'skill "pptx" loaded' in values_as_text.lower()
    else:
        pptx_skill_invoked = (
            '"skill":"pptx"' in raw_stream.replace(" ", "").lower()
            or "/pptx" in lower_evidence
            or "pptx creation, editing, and analysis" in lower_evidence
        )
        skill_md_loaded = "skills/pptx/skill.md" in lower_evidence
    progress = {
        "pptx_skill_invoked": pptx_skill_invoked,
        "skill_md_loaded": skill_md_loaded,
        "html2pptx_guide_loaded": "html2pptx.md" in tool_args_text
        or "html2pptx.md" in lower_evidence,
        "official_html2pptx_used": "html2pptx.js" in tool_args_text
        or "scripts/html2pptx.js" in lower_evidence,
        "official_thumbnail_used": "thumbnail.py" in tool_args_text
        or "scripts/thumbnail.py" in lower_evidence,
        "thumbnail_visually_inspected": any(
            term in lower_evidence for term in ("thumbnail", "overlap", "cutoff", "visual inspection")
        ),
    }
    configured_secrets = [
        value
        for name in (
            "ANTHROPIC_API_KEY", "KIMI_API_KEY", "MOONSHOT_API_KEY",
            "OPENAI_API_KEY", "OPENROUTER_API_KEY",
        )
        if (value := os.getenv(name))
    ]
    credential_scan_passed = not any(secret in raw_stream for secret in configured_secrets)
    credential_scan_passed = credential_scan_passed and not bool(
        re.search(r"sk-ant-[A-Za-z0-9_-]{12,}", raw_stream)
    )
    visual_gate = (
        len(source_visuals) >= protocol["output"]["minimum_paper_visuals"]
        and all(
            item["embedded"] and item["page"] and item["label"] and item["caption"]
            for item in source_visuals
        )
    )
    if runtime == "kimi":
        result_metadata = collect_kimi_metadata(events, run_dir)
        run_succeeded = (
            result_metadata.get("return_code") == 0
            and bool(result_metadata.get("final_response"))
            and result_metadata.get("num_tool_calls", 0) > 0
        )
        run_gate_key = "kimi_run_succeeded"
    else:
        result_metadata = collect_result_metadata(events)
        run_succeeded = bool(result_metadata) and not result_metadata.get("is_error")
        run_gate_key = "claude_run_succeeded"
    gates = {
        "source_pdf_hash_matches": paper_path.exists() and sha256(paper_path) == protocol["paper"]["pdf_sha256"],
        run_gate_key: run_succeeded,
        **progress,
        "pptx_zip_valid": zip_valid,
        "pptx_reopens": prs is not None,
        "slide_count_in_range": protocol["output"]["minimum_slides"] <= slide_count <= protocol["output"]["maximum_slides"],
        "required_sections_present": all(section_checks.values()),
        "three_source_visuals_embedded_and_documented": visual_gate,
        "thumbnail_grid_present": bool(thumbnail_candidates),
        "credential_scan_passed": credential_scan_passed,
    }
    artifacts = {}
    for path in [protocol_path, stream_path, paper_path, pptx_path, visual_manifest_path, *thumbnail_candidates]:
        if path.exists() and path.is_file():
            artifacts[str(path.relative_to(run_dir))] = {
                "sha256": sha256(path), "bytes": path.stat().st_size
            }
    return {
        "experiment_id": "2-6",
        "runtime": runtime,
        "protocol_sha256": sha256(protocol_path),
        "official_skill_receipt": json.loads((run_dir / "official_skill_receipt.json").read_text()),
        "agent_result": result_metadata,
        "slide_count": slide_count,
        "section_checks": section_checks,
        "source_visuals": source_visuals,
        "embedded_media": embedded_media,
        "thumbnail_files": [str(path.relative_to(run_dir)) for path in thumbnail_candidates],
        "gates": gates,
        "official_complete": all(gates.values()),
        "artifacts": artifacts,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("run_dir", type=Path)
    args = parser.parse_args()
    result = validate(args.run_dir)
    comparison = args.run_dir / "comparison.json"
    comparison.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    manifest = {
        "experiment_id": "2-6",
        "runtime": result["runtime"],
        "official_complete": result["official_complete"],
        "protocol_sha256": result["protocol_sha256"],
        "comparison_sha256": sha256(comparison),
        "pptx_sha256": result["artifacts"].get(
            "workspace/output/attention-is-all-you-need.pptx", {}
        ).get("sha256"),
    }
    (args.run_dir / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(json.dumps({**manifest, "gates": result["gates"]}, indent=2))
    return 0 if result["official_complete"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
